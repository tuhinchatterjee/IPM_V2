"""The pack as a document: PDF, Word, slides, and the evidence behind them.

Rendered from the SNAPSHOTS, never from a fresh calculation. That is the whole
property an export has to have: the PDF a committee is sent in the morning and
the screen somebody opens in the afternoon show the same numbers, because both
read the figures that were frozen into the pack.

Nothing here computes anything. If a figure is wrong in an export it was wrong
in the pack, which is where it can be argued with.

What is reused rather than rebuilt
-----------------------------------
    PDF and Word     `backend.reporting.writers` — one content model, two
                     containers, already written and already tested
    charts           `backend.reporting.charts` — matplotlib, rendered once as
                     PNG so every format shows the same picture
    the download log `export_records`, the platform's own table
    the workbook     openpyxl, the same library `backend/exports` uses

Slides are the one format this module writes itself, because a deck is not a
document with different margins: a committee reads slides in the room and the
pack afterwards, and the two have different jobs.

Excel formula injection
-----------------------
A cell beginning `=`, `+`, `-` or `@` is executed by Excel when the file is
opened. Committee packs carry section titles, findings and free-text
commentary, any of which can begin with a minus sign for entirely innocent
reasons — "-0.4pp on the quarter" is a sentence somebody will write. Every text
cell written by this module is prefixed so Excel treats it as text, and the
prefix is applied on write rather than by cleaning the data, because the data
is not wrong: it is a legitimate string that one program interprets as code.
"""

from __future__ import annotations

import hashlib
import io
import logging
import re
from datetime import UTC, datetime
from typing import Any

from sqlalchemy import select

from backend.models.playbook import (
    SOURCE_UI,
    PlaybookAction,
    PlaybookBlock,
    PlaybookCommittee,
    PlaybookDecision,
    PlaybookFinding,
    PlaybookSection,
    PlaybookSnapshot,
)
from backend.playbook import access, compare, readiness
from backend.playbook import snapshots as snap

logger = logging.getLogger(__name__)

#: The formats a pack can leave CreditProbe in.
FORMATS: tuple[str, ...] = ("pdf", "docx", "pptx", "xlsx")

FORMAT_LABELS: dict[str, str] = {
    "pdf": "PDF pack",
    "docx": "Word pack",
    "pptx": "Slides",
    "xlsx": "Evidence workbook",
}

MEDIA_TYPES: dict[str, str] = {
    "pdf": "application/pdf",
    "docx": ("application/vnd.openxmlformats-officedocument"
             ".wordprocessingml.document"),
    "pptx": ("application/vnd.openxmlformats-officedocument"
             ".presentationml.presentation"),
    "xlsx": ("application/vnd.openxmlformats-officedocument"
             ".spreadsheetml.sheet"),
}

#: The report writers only know three severities. A committee pack uses five,
#: so INFO and CRITICAL are mapped rather than left to raise a KeyError deep
#: inside a PDF build — which is a failure that reads as "export is broken".
SEVERITY_FOR_WRITER: dict[str, str] = {
    "CRITICAL": "HIGH", "HIGH": "HIGH", "MEDIUM": "MEDIUM",
    "LOW": "LOW", "INFO": "LOW",
}

#: Characters Excel treats as the start of a formula.
_FORMULA_START = ("=", "+", "-", "@", "\t", "\r")

#: Replaced in a filename, so a committee called `../../etc` produces a flat
#: name rather than a path.
_UNSAFE = re.compile(r"[^A-Za-z0-9._-]+")

#: Collapsed after that. A dot survives `_UNSAFE` because filenames have
#: extensions, which leaves `../../etc` as `..-..-etc` — flat and harmless, but
#: a Content-Disposition filename containing `..` is the kind of thing that
#: makes somebody look twice at a download, and there is no reason to ship it.
_DOTS = re.compile(r"\.{2,}")


def safe_cell(value: Any) -> Any:
    """A value Excel will display rather than execute.

    Numbers, dates and booleans pass through as themselves — they are not
    formulas and quoting them would turn a figure into text somebody cannot
    sum. Only strings are guarded, and only the ones that begin with a
    character Excel reads as an operator.
    """
    if value is None or isinstance(value, int | float | bool):
        return value
    text = str(value)
    if text.startswith(_FORMULA_START):
        # A leading apostrophe is Excel's own "this is text" marker. It is not
        # shown in the cell and it survives a copy-paste into another sheet,
        # which stripping the character would not.
        return f"'{text}"
    return text


def safe_filename(pack: Any, fmt: str) -> str:
    """A download name that is a name and not a path."""
    stem = _DOTS.sub(".", _UNSAFE.sub("-", f"{pack.code}-{pack.name}"))
    stem = stem[:80].strip("-.")
    return f"{stem or 'committee-pack'}.{fmt}"


# ------------------------------------------------------ the content model


def document(session: Any, pack: Any, *, include_evidence: bool = True
             ) -> dict[str, Any]:
    """The pack as `backend.reporting.writers` expects it.

    One content model, and every format renders from it, so the PDF and the
    Word file and the deck cannot disagree about what the pack says.
    """
    committee = session.get(PlaybookCommittee, int(pack.committee_id))
    sections = session.execute(
        select(PlaybookSection).where(PlaybookSection.pack_id == pack.id)
        .order_by(PlaybookSection.position)).scalars().all()
    blocks = session.execute(
        select(PlaybookBlock).where(PlaybookBlock.pack_id == pack.id)
        .order_by(PlaybookBlock.position)).scalars().all()
    figures = _figures(session, blocks)
    findings = session.execute(
        select(PlaybookFinding).where(PlaybookFinding.pack_id == pack.id)
        .order_by(PlaybookFinding.severity)).scalars().all()

    by_section: dict[int, list] = {}
    for block in blocks:
        by_section.setdefault(int(block.section_id), []).append(block)
    findings_by_section: dict[int | None, list] = {}
    for found in findings:
        findings_by_section.setdefault(found.section_id, []).append(found)

    built = []
    for section in sections:
        built.append(_section(section, by_section.get(int(section.id), []),
                              figures,
                              findings_by_section.get(int(section.id), [])))

    loose = findings_by_section.get(None) or []
    if loose:
        built.append({
            "title": "Other material findings",
            "narrative": ("Raised against this pack but not attached to a "
                          "particular section."),
            "findings": [_finding(f) for f in loose],
        })

    built.append(_decisions_section(session, pack))
    built.append(_actions_section(session, pack))
    if include_evidence:
        built.append(_basis_section(session, pack, figures))

    when = pack.approved_at or pack.updated_at or datetime.now(UTC)
    return {
        "title": safe_text(pack.name),
        "short_title": safe_text(pack.code),
        "quarter_label": safe_text(pack.period or "—"),
        "audience": safe_text(committee.name if committee is not None
                              else "Committee"),
        "prepared_by": "CreditProbe",
        "generated_at": when.strftime("%d %B %Y"),
        "purpose": safe_text(_purpose(pack, committee)),
        "classification": safe_text(_classification(pack)),
        "sections": [s for s in built if s],
    }


def _purpose(pack: Any, committee: Any) -> str:
    bits = []
    if committee is not None and str(committee.purpose or "").strip():
        bits.append(str(committee.purpose).strip())
    if pack.meeting_at is not None:
        bits.append(f"Tabled at the meeting of "
                    f"{pack.meeting_at.strftime('%d %B %Y')}.")
    if pack.as_of_date is not None:
        bits.append(f"Figures are as at "
                    f"{pack.as_of_date.strftime('%d %B %Y')}.")
    if str(pack.status) in ("APPROVED", "PUBLISHED"):
        bits.append(
            f"Approved at version {pack.approved_version}; every figure below "
            "is the one the committee was given and is not recalculated when "
            "this pack is opened.")
    else:
        bits.append(
            f"DRAFT at version {pack.version}. This is not an approved "
            "committee pack.")
    if str(pack.amendment_reason or "").strip():
        bits.append(f"Amendment: {pack.amendment_reason}")
    return " ".join(bits)


def _classification(pack: Any) -> str:
    label = str(pack.confidentiality or "CONFIDENTIAL").replace("_", " ")
    return f"{label.title()} — {pack.code}"


def _section(section: Any, blocks: list[Any], figures: dict[int, snap.Figure],
             findings: list[Any]) -> dict[str, Any]:
    """One page: its prose, its figures, its chart and its findings."""
    narrative: list[str] = []
    kpis: list[list[str]] = []
    chart = None
    table = None

    for block in blocks:
        kind = str(block.block_type)
        if kind in ("NARRATIVE", "AI_NARRATIVE", "RISK_CALLOUT",
                    "METHODOLOGY_NOTE", "DATA_QUALITY_NOTE"):
            body = str(block.body or "").strip()
            if not body:
                continue
            if kind == "AI_NARRATIVE" and not block.ai_accepted:
                # An unaccepted draft never reaches a document. Readiness
                # blocks approval on one, but a DRAFT export is legitimate and
                # must not carry words nobody has signed.
                continue
            if kind == "RISK_CALLOUT":
                body = f"Risk: {body}"
            narrative.append(body)
            continue

        figure = figures.get(int(block.snapshot_id)) if block.snapshot_id else None
        if kind == "KPI" and figure is not None:
            kpis.append([
                safe_text(figure.metric_name or figure.metric_id),
                safe_text(figure.display_value),
                safe_text(_movement_cell(figure)),
                safe_text(figure.period or "—"),
            ])
        elif kind == "CHART" and figure is not None and figure.series:
            chart = {
                "kind": str((block.config or {}).get("chart_type") or "bar"),
                "title": safe_text(block.title or figure.metric_name or ""),
                "data": list(figure.series),
            }
            chart["kind"] = ("pack_line" if chart["kind"] in ("line", "area")
                             else "pack_bar")
        elif kind == "TABLE" and figure is not None and figure.series:
            table = {
                "title": safe_text(block.title or figure.metric_name or ""),
                "columns": ["", safe_text(figure.metric_name
                                          or figure.metric_id)],
                "rows": [[safe_text(p.get("label")),
                          snap.display(p.get("value"), figure.unit,
                                       figure.decimals)]
                         for p in list(figure.series)[:40]],
            }

    built: dict[str, Any] = {"title": safe_text(section.title)}
    if str(section.purpose or "").strip():
        narrative.insert(0, str(section.purpose).strip())
    if narrative:
        # Escaped HERE rather than at each source, because every one of these
        # strings was written by a person and the PDF writer parses what it is
        # given as mini-HTML. A section legitimately called "<Finance> Review"
        # would otherwise break the build.
        built["narrative"] = safe_text("  ".join(narrative))
    if kpis:
        built["table"] = {
            "columns": ["Measure", "Value", "Movement", "Period"],
            "rows": kpis,
        }
        if table is not None:
            built["extra_table"] = table
    elif table is not None:
        built["table"] = table
    if chart is not None:
        built["chart"] = chart
    if findings:
        built["findings"] = [_finding(f) for f in findings]
    return built


def _movement_cell(figure: snap.Figure) -> str:
    moved = snap.movement(figure)
    if not moved.get("available"):
        return "—"
    arrow = {"up": "▲", "down": "▼", "flat": "="}.get(moved["direction"], "")
    tail = ""
    if moved.get("better") is False:
        tail = " (worse)"
    elif moved.get("better") is True:
        tail = " (better)"
    return f"{arrow} {moved['display']}{tail}".strip()


def _finding(row: Any) -> dict[str, str]:
    severity = SEVERITY_FOR_WRITER.get(str(row.severity), "MEDIUM")
    text = str(row.title)
    if str(row.factual_basis or "").strip():
        text = f"{text} — {row.factual_basis}"
    if str(row.status) in ("DISMISSED", "RESOLVED", "EXPLAINED"):
        text = f"{text} [{str(row.status).lower()}]"
    if str(row.response or "").strip():
        text = f"{text} Management response: {row.response}"
    return {"severity": severity, "text": safe_text(text)}


def _decisions_section(session: Any, pack: Any) -> dict[str, Any] | None:
    rows = session.execute(
        select(PlaybookDecision).where(PlaybookDecision.pack_id == pack.id)
        .order_by(PlaybookDecision.id)).scalars().all()
    if not rows:
        return None
    return {
        "title": "Decisions",
        "narrative": safe_text(
            "What this pack asks the committee to decide, and what was "
            "decided."),
        "table": {
            "columns": ["Ref", "Decision", "Recommendation", "Outcome"],
            "rows": [[safe_text(r.reference), safe_text(r.title),
                      safe_text(r.recommendation or "—"),
                      safe_text(_outcome(r))] for r in rows],
        },
    }


def _outcome(row: Any) -> str:
    from backend.playbook.actions import DECIDED

    if str(row.status) not in DECIDED:
        return str(row.status).replace("_", " ").title()
    when = row.decided_at.strftime("%d %b %Y") if row.decided_at else ""
    said = str(row.status).replace("_", " ").title()
    return f"{said} {when}".strip()


def _actions_section(session: Any, pack: Any) -> dict[str, Any] | None:
    """Every open action the COMMITTEE is carrying, not only this pack's.

    An action agreed three meetings ago is still the committee's business,
    and a pack that showed only its own would let one quietly disappear.
    """
    from backend.playbook.actions import CLOSED

    rows = session.execute(
        select(PlaybookAction).where(
            PlaybookAction.committee_id == pack.committee_id,
            PlaybookAction.status.notin_(tuple(CLOSED) + ("DRAFT",)))
        .order_by(PlaybookAction.due_date.asc().nullslast())).scalars().all()
    if not rows:
        return None
    return {
        "title": "Action log",
        "narrative": ("Every action this committee is carrying. Progress is "
                      "read from the Project Planner where an action has been "
                      "linked to one."),
        "table": {
            "columns": ["Ref", "Action", "Owner", "Due", "Status", "Update"],
            "rows": [[
                safe_text(r.reference), safe_text(r.description[:160]),
                _person(session, r.owner_id),
                r.due_date.strftime("%d %b %Y") if r.due_date else "—",
                str(r.status).replace("_", " ").title(),
                safe_text((r.latest_update or "No update since it was "
                           "raised.")[:160]),
            ] for r in rows],
        },
    }


def _basis_section(session: Any, pack: Any,
                   figures: dict[int, snap.Figure]) -> dict[str, Any]:
    """Where every figure came from, and every figure that has no value.

    The appendix a challenge session opens first. Listing the absent figures
    with their reasons is the point: a pack that simply omits them looks
    complete, and the reader has no way to know what they were not shown.
    """
    rows = []
    missing = []
    seen: set[str] = set()
    for figure in figures.values():
        if figure.metric_id in seen:
            continue
        seen.add(figure.metric_id)
        rows.append([
            safe_text(figure.metric_name or figure.metric_id),
            safe_text(figure.metric_id),
            safe_text(figure.period or "—"),
            safe_text(figure.dataset or "—"),
            safe_text(figure.metric_version or "—"),
            safe_text((figure.formula_hash or "")[:12] or "—"),
        ])
        if not figure.available:
            missing.append({
                "severity": "MEDIUM",
                "text": safe_text(
                    f"{figure.metric_name or figure.metric_id} has no value "
                    f"for {figure.period or 'this period'} "
                    f"({figure.availability}): {figure.unavailable_reason}"),
            })

    state = readiness.assess(session, pack)
    narrative = [
        "Every figure in this pack was calculated by CreditProbe's governed "
        "metric catalogue at the version and period shown, and stored with "
        "the pack. Opening this pack does not recalculate anything, so the "
        "figures here are the ones the committee was given.",
    ]
    if str(pack.status) in ("APPROVED", "PUBLISHED"):
        narrative.append(
            f"Approved at version {pack.approved_version} on "
            f"{pack.approved_at.strftime('%d %B %Y') if pack.approved_at else 'the date recorded'}.")
    else:
        narrative.append(
            f"This is a draft at version {pack.version}. Readiness at the "
            f"time of export: {state.percent}% ({state.state}), with "
            f"{len(state.blocking)} blocking item"
            f"{'s' if len(state.blocking) != 1 else ''}.")

    built: dict[str, Any] = {
        "title": "Basis of preparation",
        "narrative": safe_text(" ".join(narrative)),
    }
    if rows:
        built["table"] = {
            "columns": ["Measure", "Metric", "Period", "Source", "Version",
                        "Formula"],
            "rows": rows,
        }
    if missing:
        built["findings"] = missing
    return built


def safe_text(value: Any) -> str:
    """Text for a document, with nothing that could be read as markup.

    The PDF writer passes strings through reportlab's mini-HTML parser, so a
    section title containing a `<` would break the build — and the title is
    written by a person, who may legitimately call a section "<Finance>
    Review". Escaped rather than stripped, so the title still READS the way
    they wrote it.
    """
    text = str(value if value is not None else "")
    return (text.replace("&", "&amp;").replace("<", "&lt;")
            .replace(">", "&gt;"))


def _person(session: Any, user_id: int | None) -> str:
    if user_id is None:
        return "—"
    from backend.db.models import User

    row = session.get(User, int(user_id))
    if row is None:
        return "—"
    name = " ".join(p for p in (str(row.first_name or ""),
                                str(row.last_name or "")) if p).strip()
    return safe_text(name or str(row.username))


def _figures(session: Any, blocks: list[Any]) -> dict[int, snap.Figure]:
    ids = [int(b.snapshot_id) for b in blocks if b.snapshot_id is not None]
    if not ids:
        return {}
    rows = session.execute(
        select(PlaybookSnapshot)
        .where(PlaybookSnapshot.id.in_(ids))).scalars().all()
    return {int(r.id): snap.from_row(r) for r in rows}


# ------------------------------------------------------------- the formats


def render(session: Any, pack: Any, fmt: str) -> tuple[bytes, str]:
    """One pack, one format. Returns (bytes, media type)."""
    wanted = str(fmt or "").lower()
    if wanted not in FORMATS:
        from backend.playbook.service import InvalidPlaybook

        raise InvalidPlaybook(
            f"'{fmt}' is not a format a committee pack can be exported in. "
            f"One of: {', '.join(FORMATS)}.")

    if wanted == "xlsx":
        return _workbook(session, pack), MEDIA_TYPES["xlsx"]

    built = document(session, pack)
    if wanted == "pptx":
        return _slides(session, pack, built), MEDIA_TYPES["pptx"]

    from backend.reporting import writers

    if wanted == "pdf":
        return writers.write_pdf(built), MEDIA_TYPES["pdf"]
    return writers.write_docx(built), MEDIA_TYPES["docx"]


def _slides(session: Any, pack: Any, built: dict[str, Any]) -> bytes:
    """The deck a chair presents from.

    Not the document with different margins. A slide carries the figures and
    the material findings and leaves the working to the pack, because a
    committee reads slides in the room and the pack afterwards, and putting
    the basis of preparation on a slide is how nobody reads either.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    navy = RGBColor(0x0B, 0x24, 0x36)
    muted = RGBColor(0x6C, 0x7A, 0x8C)

    cover = deck.slides.add_slide(deck.slide_layouts[6])
    _text(cover, built["title"], Inches(0.8), Inches(2.4), Inches(11.7),
          Inches(1.2), size=34, bold=True, colour=navy)
    _text(cover, f"{built['audience']} · {built['quarter_label']}",
          Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6), size=18,
          colour=muted)
    _text(cover, built["classification"], Inches(0.8), Inches(6.6),
          Inches(11.7), Inches(0.4), size=11, colour=muted)

    for section in built["sections"]:
        if section["title"] == "Basis of preparation":
            # Deliberately left off the deck. It belongs in the pack, and a
            # slide of formula hashes is a slide nobody reads.
            continue
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        _text(slide, section["title"], Inches(0.6), Inches(0.4),
              Inches(12.1), Inches(0.8), size=26, bold=True, colour=navy)

        top = Inches(1.4)
        table = section.get("table")
        if table and table.get("rows"):
            rows = table["rows"][:8]
            shape = slide.shapes.add_table(
                len(rows) + 1, len(table["columns"]), Inches(0.6), top,
                Inches(12.1), Inches(0.4 * (len(rows) + 1))).table
            for column, name in enumerate(table["columns"]):
                cell = shape.cell(0, column)
                cell.text = str(name)
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
            for index, row in enumerate(rows, 1):
                for column, value in enumerate(row):
                    cell = shape.cell(index, column)
                    cell.text = _plain(value)
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            top = Inches(1.4 + 0.42 * (len(rows) + 1))

        said = str(section.get("narrative") or "")
        if said:
            _text(slide, _plain(said)[:600], Inches(0.6), top, Inches(12.1),
                  Inches(1.6), size=13)
            top = top + Inches(1.8)

        findings = section.get("findings") or []
        if findings:
            lines = "\n".join(f"[{f['severity']}] {_plain(f['text'])[:180]}"
                              for f in findings[:4])
            _text(slide, lines, Inches(0.6), top, Inches(12.1), Inches(1.8),
                  size=12, colour=RGBColor(0xC0, 0x29, 0x2E))

    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _text(slide: Any, text: str, left: Any, top: Any, width: Any, height: Any,
          *, size: int = 14, bold: bool = False, colour: Any = None) -> None:
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = text
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            if colour is not None:
                run.font.color.rgb = colour


def _plain(text: Any) -> str:
    """Undo the document escaping for a format that does not parse markup."""
    return (str(text).replace("&lt;", "<").replace("&gt;", ">")
            .replace("&amp;", "&"))


def _workbook(session: Any, pack: Any) -> bytes:
    """The evidence behind the pack, one sheet per kind of thing.

    What somebody opens when they want to check a figure rather than read a
    narrative: every snapshot with its formula, period, numerator,
    denominator, dataset version and run id.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    book = Workbook()
    head = Font(bold=True, color="FFFFFF")
    fill = PatternFill("solid", fgColor="0B2436")

    def sheet(title: str, columns: list[str], rows: list[list[Any]]) -> None:
        page = book.create_sheet(title[:31])
        page.append([safe_cell(c) for c in columns])
        for cell in page[1]:
            cell.font = head
            cell.fill = fill
            cell.alignment = Alignment(vertical="center")
        for row in rows:
            page.append([safe_cell(v) for v in row])
        for index, name in enumerate(columns, 1):
            width = max(len(str(name)) + 2,
                        *(len(str(r[index - 1])[:60]) + 2 for r in rows)) \
                if rows else len(str(name)) + 2
            page.column_dimensions[
                page.cell(row=1, column=index).column_letter].width = min(
                    60, max(10, width))
        page.freeze_panes = "A2"

    book.remove(book.active)
    committee = session.get(PlaybookCommittee, int(pack.committee_id))
    state = readiness.assess(session, pack)

    sheet("PACK", ["Property", "Value"], [
        ["Code", pack.code],
        ["Name", pack.name],
        ["Committee", committee.name if committee is not None else ""],
        ["Reporting period", pack.period],
        ["Comparison period", pack.comparison_period],
        ["Meeting", pack.meeting_at.isoformat() if pack.meeting_at else ""],
        ["As at", pack.as_of_date.isoformat() if pack.as_of_date else ""],
        ["Status", pack.status],
        ["Working version", pack.version],
        ["Approved version", pack.approved_version],
        ["Approved at", pack.approved_at.isoformat() if pack.approved_at else ""],
        ["Confidentiality", pack.confidentiality],
        ["Readiness at export", f"{state.percent}% {state.state}"],
        ["Blocking items", len(state.blocking)],
        ["Exported", datetime.now(UTC).isoformat()],
    ])

    blocks = session.execute(
        select(PlaybookBlock).where(PlaybookBlock.pack_id == pack.id)
        .order_by(PlaybookBlock.position)).scalars().all()
    figures = _figures(session, blocks)
    rows = []
    for figure in figures.values():
        rows.append([
            figure.metric_id, figure.metric_name, figure.period,
            figure.display_value, figure.value, figure.availability,
            figure.unavailable_reason, figure.numerator, figure.denominator,
            figure.rows_considered, figure.unit, figure.metric_version,
            figure.formula_hash, figure.dataset, figure.dataset_version,
            figure.run_id, figure.verification_state,
            "governed" if figure.governed else "user-defined",
        ])
    sheet("FIGURES", [
        "Metric", "Name", "Period", "Shown as", "Value", "Availability",
        "Why no value", "Numerator", "Denominator", "Rows considered", "Unit",
        "Metric version", "Formula hash", "Source", "Source version", "Run id",
        "Verification", "Origin",
    ], rows)

    findings = session.execute(
        select(PlaybookFinding)
        .where(PlaybookFinding.pack_id == pack.id)).scalars().all()
    sheet("FINDINGS", [
        "Severity", "Type", "Title", "Status", "Rule", "Threshold",
        "Factual basis", "Response", "Dismissed because",
    ], [[
        f.severity, f.finding_type, f.title, f.status, f.rule_key,
        (f.rule_detail or {}).get("rule", {}).get("threshold"),
        f.factual_basis, f.response, f.dismissed_reason,
    ] for f in findings])

    decisions = session.execute(
        select(PlaybookDecision)
        .where(PlaybookDecision.pack_id == pack.id)).scalars().all()
    sheet("DECISIONS", [
        "Ref", "Title", "Question", "Recommendation", "Status", "Decided",
        "Decision", "Conditions",
    ], [[
        d.reference, d.title, d.question, d.recommendation, d.status,
        d.decided_at.isoformat() if d.decided_at else "",
        d.decision_text, d.conditions,
    ] for d in decisions])

    from backend.playbook.actions import CLOSED

    action_rows = session.execute(
        select(PlaybookAction).where(
            PlaybookAction.committee_id == pack.committee_id,
            PlaybookAction.status.notin_(tuple(CLOSED) + ("DRAFT",)))
    ).scalars().all()
    sheet("ACTIONS", [
        "Ref", "Description", "Owner", "Due", "Priority", "Status",
        "Latest update", "Planner task",
    ], [[
        a.reference, a.description, _plain(_person(session, a.owner_id)),
        a.due_date.isoformat() if a.due_date else "", a.priority, a.status,
        a.latest_update, a.planner_task_id,
    ] for a in action_rows])

    changed = compare.against(session, pack, compare._previous(session, pack))
    sheet("SINCE LAST", [
        "Metric", "Name", "Kind", "Now", "Then", "Change", "Caveat",
    ], [[
        d.metric_id, d.name, d.kind, d.now_display, d.then_display,
        d.change_display, d.caveat,
    ] for d in changed.material])

    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------- the door


def export(session: Any, pack_id: int, principal: Any, *, fmt: str = "pdf",
           source: str = SOURCE_UI) -> dict[str, Any]:
    """Produce one export and record that it happened.

    The authorisation is the pack's own: anybody who may READ the pack may
    export it, because a screenshot of the screen is the same disclosure and
    a permission that stops one and not the other is theatre. What matters —
    and what is recorded — is WHO took a copy and WHEN.
    """
    from backend.models.platform import ExportRecord
    from backend.playbook.service import record

    pack, grant = access.readable_pack(session, pack_id, principal, source)
    data, media = render(session, pack, fmt)
    filename = safe_filename(pack, fmt.lower())
    checksum = hashlib.sha256(data).hexdigest()

    # The row is deliberately wide, for the reason `ExportRecord` records: six
    # months later the question is not "did an export happen" but "exactly
    # which figures, from which data version, did that file contain". Every
    # snapshot in the file is named here, so answering it never needs the file
    # regenerated — which would answer with today's figures anyway.
    blocks = session.execute(
        select(PlaybookBlock).where(
            PlaybookBlock.pack_id == pack.id)).scalars().all()
    figures = _figures(session, blocks)
    session.add(ExportRecord(
        kind=f"playbook_{fmt.lower()}",
        object_type="playbook_pack", object_id=str(pack.id),
        user_id=grant.user_id,
        role=str(getattr(getattr(principal, "role", ""), "value",
                         getattr(principal, "role", "")) or ""),
        status="allowed",
        authorization=(f"{grant.access.lower()}:committee_"
                       f"{pack.committee_id}"
                       f"{':administrative' if grant.administrative else ''}"
                       )[:64],
        filename=filename, content_hash=checksum, size_bytes=len(data),
        row_count=len(figures),
        datasets=sorted({f.dataset for f in figures.values() if f.dataset}),
        detail={
            "pack_code": str(pack.code),
            "pack_status": str(pack.status),
            "pack_version": int(pack.version),
            "approved_version": pack.approved_version,
            "period": str(pack.period),
            "confidentiality": str(pack.confidentiality),
            "source": grant.source,
            "figures": [
                {"metric_id": f.metric_id, "period": f.period,
                 "display_value": f.display_value,
                 "availability": f.availability,
                 "metric_version": f.metric_version,
                 "formula_hash": f.formula_hash,
                 "dataset_version": f.dataset_version, "run_id": f.run_id}
                for f in figures.values()],
        }))
    session.flush()

    record(session, entity_type="pack", action="exported", pack=pack,
           entity_id=int(pack.id), entity_ref=str(pack.code),
           narrative=(f"Exported as {FORMAT_LABELS.get(fmt.lower(), fmt)} "
                      f"({len(data):,} bytes, sha256 {checksum[:12]})."),
           grant=grant)
    return {
        "filename": filename, "media_type": media, "bytes": data,
        "size": len(data), "checksum": checksum,
        "format": fmt.lower(), "label": FORMAT_LABELS.get(fmt.lower(), fmt),
    }


def formats() -> list[dict[str, str]]:
    """What a download button may offer, and what each one is for."""
    return [
        {"format": "pdf", "label": FORMAT_LABELS["pdf"],
         "purpose": "The pack as it is circulated and filed."},
        {"format": "docx", "label": FORMAT_LABELS["docx"],
         "purpose": "The same pack, editable, for a secretariat that "
                    "assembles a board bundle."},
        {"format": "pptx", "label": FORMAT_LABELS["pptx"],
         "purpose": "The deck a chair presents from. Figures and material "
                    "findings; the working stays in the pack."},
        {"format": "xlsx", "label": FORMAT_LABELS["xlsx"],
         "purpose": "Every figure with its formula, period, numerator, "
                    "denominator, source version and run id."},
    ]


__all__ = [
    "FORMATS", "FORMAT_LABELS", "MEDIA_TYPES", "SEVERITY_FOR_WRITER",
    "document", "export", "formats", "render", "safe_cell", "safe_filename",
    "safe_text",
]
