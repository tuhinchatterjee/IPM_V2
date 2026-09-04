"""Reading somebody's existing pack in, without pretending CreditProbe made it.

A committee that has been meeting for eleven years has a pack already. Getting
started must not mean rebuilding it from nothing, so this module takes the file
they have and lays out a draft from it.

Everything it produces is labelled
-----------------------------------
A number lifted out of a spreadsheet is NOT a governed CreditProbe figure, and
the difference is the whole product. Every block this module creates carries an
`import_class`:

    IMPORTED_TEXT           prose, taken as-is
    IMPORTED_IMAGE          a picture nobody can interrogate
    UNMAPPED_TABLE          a table whose numbers came from the file
    MAPPED_GOVERNED_METRIC  a person has since pointed this at a real metric
    SUPPORTING_DOCUMENT     the file itself, kept as evidence

The first three read on screen and in every export as "from the uploaded
document". Only the fourth is a figure the platform stands behind, and it only
becomes one when a person maps it.

Uploaded content is never trusted
----------------------------------
Three separate things, all of them real:

**Prompt injection.** A paragraph in somebody's pack saying "ignore your
instructions" reaches nothing that reads instructions: `backend.playbook.
narrative` builds its prompt from governed figures and the section's own
configuration, and imported text is not in it. That is the defence — not a
filter on the words, which would be a filter somebody eventually gets past.

**The file itself.** Size, type and shape are checked before anything is
parsed, a zip is refused if its declared contents are larger than the limit
(a zip bomb declares its size honestly in the central directory), and the
filename is never used as a path.

**Excel formulas.** A cell beginning `=` is executed on open. Imported cell
VALUES are read, never the formulas, and anything written back out goes
through `export.safe_cell`.
"""

from __future__ import annotations

import hashlib
import logging
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from backend.config import settings
from backend.models.playbook import (
    SOURCE_IMPORT,
    PlaybookSource,
)
from backend.playbook import access, service
from backend.playbook.access import CONTRIBUTOR
from backend.playbook.service import InvalidPlaybook

logger = logging.getLogger(__name__)

#: The largest file this will accept. A committee pack is a document, not a
#: data extract; forty megabytes of PowerPoint is somebody uploading the wrong
#: thing, and reading it would tie up a worker for minutes.
MAX_BYTES = 40 * 1024 * 1024

#: The largest total the contents of a zip-based file (docx, xlsx, pptx) may
#: declare. Checked from the central directory BEFORE anything is decompressed,
#: which is what makes it a defence against a zip bomb rather than a report of
#: one: a 40 KB file that declares 8 GB is refused without being read.
MAX_UNPACKED = 400 * 1024 * 1024

#: And the largest ratio. A legitimate Office file compresses maybe 20:1; a
#: bomb is 1000:1 and upwards.
MAX_RATIO = 120

#: What can actually be read. Anything else is kept as a supporting document
#: rather than refused, because a committee's PDF appendix is legitimate
#: evidence even though nothing here can extract a table from it.
READABLE: dict[str, str] = {
    ".docx": "word",
    ".xlsx": "excel",
    ".pptx": "slides",
}

SUPPORTING: frozenset[str] = frozenset({
    ".pdf", ".doc", ".xls", ".ppt", ".csv", ".txt", ".md", ".png", ".jpg",
    ".jpeg", ".gif", ".svg",
})

#: Content types a browser sends for the files above. Checked as a HINT, never
#: as the authority: a client controls its own Content-Type header, so the
#: extension and the file's own magic bytes are what actually decide.
MEDIA_HINTS: dict[str, tuple[str, ...]] = {
    ".docx": ("application/vnd.openxmlformats-officedocument"
              ".wordprocessingml.document",),
    ".xlsx": ("application/vnd.openxmlformats-officedocument"
              ".spreadsheetml.sheet",),
    ".pptx": ("application/vnd.openxmlformats-officedocument"
              ".presentationml.presentation",),
    ".pdf": ("application/pdf",),
}

#: A zip-based Office file starts with these. Checked so a `.docx` that is
#: really something else is refused before a parser sees it.
_ZIP_MAGIC = (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")

#: Stripped from a stored filename. The name is used for display and for a
#: file on disk, and both want the same restriction.
_UNSAFE = re.compile(r"[^A-Za-z0-9._-]+")
_DOTS = re.compile(r"\.{2,}")

#: Paragraphs shorter than this are headings or stray fragments rather than
#: commentary. Kept out of narrative blocks so an imported pack does not
#: produce forty one-word sections.
MIN_PARAGRAPH = 40


class Unreadable(InvalidPlaybook):
    """The file cannot be read, and the reason is worth saying."""


@dataclass
class Imported:
    """What one import produced, and what it could not do."""

    source_id: int | None = None
    filename: str = ""
    kind: str = ""
    sections: int = 0
    blocks: int = 0
    tables: int = 0
    paragraphs: int = 0
    warnings: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "source_id": self.source_id, "filename": self.filename,
            "kind": self.kind, "sections": self.sections,
            "blocks": self.blocks, "tables": self.tables,
            "paragraphs": self.paragraphs, "warnings": list(self.warnings),
            "summary": self.summary,
        }

    @property
    def summary(self) -> str:
        if not self.blocks:
            return (f"{self.filename} was kept as a supporting document. "
                    "Nothing in it could be read into pack content.")
        return (
            f"{self.filename} produced {self.sections} section"
            f"{'s' if self.sections != 1 else ''} and {self.blocks} block"
            f"{'s' if self.blocks != 1 else ''} — {self.paragraphs} of text "
            f"and {self.tables} table{'s' if self.tables != 1 else ''}. Every "
            "one is marked as coming from the uploaded document, not as a "
            "CreditProbe figure, until somebody maps it to a governed metric.")


def safe_name(filename: str) -> str:
    """A filename that is a name.

    A browser sends whatever the operating system had, including `../` on a
    client that wanted to. Used for display AND for the stored path, so it is
    cleaned once here rather than at each use.
    """
    stem = _DOTS.sub(".", _UNSAFE.sub("-", str(filename or ""))).strip("-.")
    return stem[:120] or "uploaded-file"


def inspect(data: bytes, filename: str,
            content_type: str = "") -> tuple[str, str, list[str]]:
    """What this file is, before anything parses it.

    Returns `(extension, kind, warnings)`. Raises `Unreadable` where the file
    should not be opened at all.

    The order matters: size, then extension, then magic bytes, then the zip's
    declared contents. Each check is cheaper than the one after it, and the
    expensive ones never run on a file the cheap ones would have refused.
    """
    warnings: list[str] = []
    if not data:
        raise Unreadable("That file is empty.")
    if len(data) > MAX_BYTES:
        raise Unreadable(
            f"That file is {len(data) / 1024 / 1024:.1f} MB and the limit is "
            f"{MAX_BYTES // 1024 // 1024} MB. A committee pack is a document; "
            "a file this size is usually a data extract, which belongs in "
            "Data Builder.")

    suffix = Path(safe_name(filename)).suffix.lower()
    if not suffix:
        raise Unreadable(
            "That file has no extension, so there is no way to tell what it "
            "is. Rename it and upload it again.")

    known = MEDIA_HINTS.get(suffix)
    if known and content_type and content_type not in known:
        # A hint, not the authority. Recorded rather than acted on: the
        # browser's Content-Type is set by the client and proves nothing, and
        # a legitimate upload from an unusual client would be refused by a
        # rule that trusted it.
        warnings.append(
            f"The browser called this a {content_type}, and the extension "
            f"says {suffix}. The extension and the file's own contents were "
            "used.")

    if suffix in READABLE:
        if not data.startswith(_ZIP_MAGIC):
            raise Unreadable(
                f"That file is named {suffix} but is not one — a {suffix} is a "
                "zip archive and this does not begin like one. Check it "
                "opens on your own machine first.")
        warnings.extend(_check_archive(data, suffix))
        return suffix, READABLE[suffix], warnings

    if suffix in SUPPORTING:
        return suffix, "supporting", warnings

    raise Unreadable(
        f"CreditProbe cannot read a {suffix} file. It reads "
        f"{', '.join(sorted(READABLE))} into pack content, and keeps "
        f"{', '.join(sorted(SUPPORTING))} as supporting evidence.")


def _check_archive(data: bytes, suffix: str) -> list[str]:
    """Refuse a zip whose DECLARED contents are absurd.

    Read from the central directory, so nothing is decompressed to find out.
    A zip bomb is honest about its uncompressed size — that is how the format
    works — which is what makes this check cheap and effective. A file that
    lies about it fails later at the parser, on a bounded read.
    """
    import io

    warnings: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = archive.infolist()
            declared = sum(int(e.file_size) for e in entries)
    except zipfile.BadZipFile as e:
        raise Unreadable(
            f"That {suffix} could not be opened as an archive: {e}") from e

    if declared > MAX_UNPACKED:
        raise Unreadable(
            f"That file declares {declared / 1024 / 1024:.0f} MB of contents "
            f"inside {len(data) / 1024:.0f} KB. CreditProbe will not unpack "
            "it.")
    if len(data) and declared / max(1, len(data)) > MAX_RATIO:
        raise Unreadable(
            f"That file compresses {declared / max(1, len(data)):.0f} to one, "
            "which is far outside what a document does. CreditProbe will not "
            "unpack it.")
    if len(entries) > 5_000:
        warnings.append(
            f"The file contains {len(entries)} parts, which is unusual for a "
            "document. Only the text and tables were read.")
    return warnings


# ------------------------------------------------------------- the readers


def read(data: bytes, kind: str) -> tuple[list[dict[str, Any]], list[str]]:
    """Turn a file into candidate sections, with what could not be read.

    Every reader returns the same shape:
    `[{"title": ..., "paragraphs": [...], "tables": [{"columns", "rows"}]}]`
    so the caller does not branch on format.
    """
    if kind == "word":
        return _read_word(data)
    if kind == "excel":
        return _read_excel(data)
    if kind == "slides":
        return _read_slides(data)
    return [], ["This file is kept as supporting evidence; nothing was read "
                "out of it."]


def _read_word(data: bytes) -> tuple[list[dict[str, Any]], list[str]]:
    """Headings become sections; paragraphs and tables become their content."""
    import io

    from docx import Document

    warnings: list[str] = []
    document = Document(io.BytesIO(data))
    sections: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None

    for paragraph in document.paragraphs:
        text = str(paragraph.text or "").strip()
        if not text:
            continue
        style = str(getattr(paragraph.style, "name", "") or "")
        if style.startswith("Heading") or style == "Title":
            current = {"title": text[:240], "paragraphs": [], "tables": []}
            sections.append(current)
            continue
        if current is None:
            current = {"title": "Imported content", "paragraphs": [],
                       "tables": []}
            sections.append(current)
        if len(text) >= MIN_PARAGRAPH:
            current["paragraphs"].append(text)

    for table in document.tables:
        rows = [[str(c.text or "").strip() for c in r.cells]
                for r in table.rows]
        if len(rows) < 2:
            warnings.append(
                "A table with no body rows was skipped — a header on its own "
                "is not data.")
            continue
        if current is None:
            current = {"title": "Imported tables", "paragraphs": [],
                       "tables": []}
            sections.append(current)
        current["tables"].append({"columns": rows[0], "rows": rows[1:]})

    if not sections:
        warnings.append(
            "The document had no headings, paragraphs or tables this could "
            "read. It is kept as a supporting document.")
    return sections, warnings


def _read_excel(data: bytes) -> tuple[list[dict[str, Any]], list[str]]:
    """One section per sheet, one table per sheet.

    `data_only=True` reads the cached VALUES rather than the formulas, which
    is both what a committee pack means and what keeps a formula out of the
    imported content entirely.
    """
    import io

    from openpyxl import load_workbook

    warnings: list[str] = []
    book = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sections: list[dict[str, Any]] = []

    for name in book.sheetnames[:30]:
        page = book[name]
        rows = []
        for row in page.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in values):
                rows.append(values)
            if len(rows) > 500:
                warnings.append(
                    f"Sheet “{name}” has more than 500 rows; the first 500 "
                    "were read. A pack table longer than that belongs in Data "
                    "Builder as a dataset.")
                break
        if len(rows) < 2:
            continue
        sections.append({
            "title": str(name)[:240], "paragraphs": [],
            "tables": [{"columns": rows[0], "rows": rows[1:]}],
        })

    if len(book.sheetnames) > 30:
        warnings.append(
            f"The workbook has {len(book.sheetnames)} sheets; the first 30 "
            "were read.")
    if not sections:
        warnings.append("No sheet in that workbook had readable rows.")
    return sections, warnings


def _read_slides(data: bytes) -> tuple[list[dict[str, Any]], list[str]]:
    """One section per slide. Pictures are noted, not extracted."""
    import io

    from pptx import Presentation

    warnings: list[str] = []
    deck = Presentation(io.BytesIO(data))
    sections: list[dict[str, Any]] = []
    pictures = 0

    for index, slide in enumerate(deck.slides, 1):
        title = ""
        paragraphs: list[str] = []
        tables: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                pictures += 1
                continue
            if getattr(shape, "has_table", False):
                rows = [[str(c.text or "").strip() for c in r.cells]
                        for r in shape.table.rows]
                if len(rows) >= 2:
                    tables.append({"columns": rows[0], "rows": rows[1:]})
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = str(shape.text_frame.text or "").strip()
            if not text:
                continue
            if not title:
                title = text.splitlines()[0][:240]
                remainder = "\n".join(text.splitlines()[1:]).strip()
                if len(remainder) >= MIN_PARAGRAPH:
                    paragraphs.append(remainder)
                continue
            if len(text) >= MIN_PARAGRAPH:
                paragraphs.append(text)
        if title or paragraphs or tables:
            sections.append({"title": title or f"Slide {index}",
                             "paragraphs": paragraphs, "tables": tables})

    if pictures:
        warnings.append(
            f"{pictures} picture{'s were' if pictures != 1 else ' was'} left "
            "behind. CreditProbe imports text and tables; an image is not "
            "something a reader can interrogate, and importing one as pack "
            "content would put a number on a page with no way to check it.")
    if not sections:
        warnings.append("No slide in that deck had readable text or tables.")
    return sections, warnings


# ----------------------------------------------------------------- the door


def import_pack(session: Any, pack_id: int, principal: Any, *, data: bytes,
                filename: str, content_type: str = "",
                as_content: bool = True,
                source: str = SOURCE_IMPORT) -> Imported:
    """Read a file into a draft pack, labelling everything it produces.

    `as_content=False` keeps the file as supporting evidence and creates no
    blocks, which is what somebody does with a PDF appendix.

    `source` is the door this came through. It defaults to IMPORT because the
    router is the only caller today, and it is a PARAMETER rather than a
    constant so the refusal below is reachable: a guard no caller can trigger
    is a guard nobody has tested, and it reads exactly like a working one.
    """
    pack, grant = access.writable_pack(
        session, pack_id, principal, CONTRIBUTOR,
        "import a document into this pack", source)
    access.refuse_ai(grant, "import_document")

    suffix, kind, warnings = inspect(data, filename, content_type)
    name = safe_name(filename)
    checksum = hashlib.sha256(data).hexdigest()
    stored = _store(data, pack, name, checksum)

    row = PlaybookSource(
        pack_id=int(pack.id), kind="DOCUMENT", label=name,
        reference=f"upload:{checksum[:16]}",
        detail={"kind": kind, "suffix": suffix,
                "declared_content_type": content_type},
        filename=name, content_type=content_type[:120],
        byte_size=len(data), stored_path=str(stored), checksum=checksum,
        import_class="SUPPORTING_DOCUMENT", warnings=list(warnings),
        uploaded_by=grant.user_id)
    session.add(row)
    session.flush()

    outcome = Imported(source_id=int(row.id), filename=name, kind=kind,
                       warnings=list(warnings))
    if as_content and kind != "supporting":
        try:
            found, more = read(data, kind)
        except Unreadable:
            raise
        except Exception as e:  # noqa: BLE001 - reported against the upload
            logger.warning("could not read %s", name, exc_info=True)
            outcome.warnings.append(
                f"CreditProbe could not read the contents of that file, so it "
                f"is kept as a supporting document only: {e}")
            found, more = [], []
        outcome.warnings.extend(more)
        _lay_out(session, pack, principal, found, outcome, int(row.id))

    row.warnings = list(outcome.warnings)
    session.flush()
    service.record(
        session, entity_type="pack", action="imported", pack=pack,
        entity_id=int(row.id), entity_ref=name,
        narrative=outcome.summary, grant=grant)
    return outcome


def _lay_out(session: Any, pack: Any, principal: Any,
             found: list[dict[str, Any]], outcome: Imported,
             source_id: int) -> None:
    """Turn read content into sections and labelled blocks."""
    for spec in found:
        section = service.create_section(
            session, int(pack.id), principal,
            title=str(spec.get("title") or "Imported content"),
            purpose="Imported from an uploaded document.",
            required=False, source=SOURCE_IMPORT)
        outcome.sections += 1

        for text in list(spec.get("paragraphs") or []):
            service.create_block(
                session, int(section["id"]), principal,
                block_type="NARRATIVE", body=str(text),
                statement_kind="NOT_RECORDED",
                import_class="IMPORTED_TEXT",
                config={"import_source_id": source_id},
                source=SOURCE_IMPORT)
            outcome.blocks += 1
            outcome.paragraphs += 1

        for table in list(spec.get("tables") or []):
            service.create_block(
                session, int(section["id"]), principal,
                block_type="TABLE",
                title=str(spec.get("title") or "Imported table")[:240],
                # No metric_id, and UNMAPPED_TABLE is why: this table's
                # numbers came out of a file and CreditProbe did not calculate
                # them. A person maps it to a governed metric, or it stays
                # labelled as theirs for as long as the pack exists.
                import_class="UNMAPPED_TABLE",
                config={
                    "imported": True,
                    "import_source_id": source_id,
                    "columns": list(table.get("columns") or []),
                    "rows": list(table.get("rows") or [])[:200],
                },
                source=SOURCE_IMPORT)
            outcome.blocks += 1
            outcome.tables += 1


def _store(data: bytes, pack: Any, name: str, checksum: str) -> Path:
    """Write the bytes somewhere they can be fetched again.

    Under a directory named by the pack id and a hash, never by the uploaded
    filename: a name is what an attacker controls, and joining one onto a path
    is the whole of a path-traversal write.
    """
    root = Path(settings.upload_dir) / "playbook" / str(int(pack.id))
    root.mkdir(parents=True, exist_ok=True)
    suffix = Path(name).suffix.lower()[:12]
    target = root / f"{checksum[:32]}{suffix}"
    target.write_bytes(data)
    return target


def sources(session: Any, pack_id: int, principal: Any) -> list[dict[str, Any]]:
    """Every document attached to this pack."""
    from sqlalchemy import select

    pack, _ = access.readable_pack(session, pack_id, principal)
    rows = session.execute(
        select(PlaybookSource).where(PlaybookSource.pack_id == pack.id)
        .order_by(PlaybookSource.id)).scalars().all()
    return [{
        "id": int(r.id), "kind": str(r.kind), "label": str(r.label),
        "filename": str(r.filename), "content_type": str(r.content_type),
        "byte_size": int(r.byte_size), "checksum": str(r.checksum),
        "import_class": str(r.import_class),
        "warnings": list(r.warnings or []),
        "uploaded_by": r.uploaded_by,
        "created_at": r.created_at.isoformat() if r.created_at else None,
    } for r in rows]


def map_to_metric(session: Any, block_id: int, principal: Any, *,
                  metric_id: str) -> dict[str, Any]:
    """Point an imported table at a governed metric.

    The moment an imported figure becomes a CreditProbe figure — and the only
    way it does. The block is reclassified and will be calculated on the next
    generation like any other, so what a reader sees afterwards is the
    platform's number rather than the file's.
    """
    from backend.metrics import service as metrics

    block, section, pack, grant = access.visible_block(
        session, block_id, principal)
    access.assert_editable(pack)
    access.may_edit_section(session, section, grant,
                            "map this table to a metric")
    if not str(block.import_class or ""):
        raise InvalidPlaybook(
            "That block was not imported, so it already shows a governed "
            "figure or is a person's own words.")
    try:
        metrics.resolve(metric_id, user_id=grant.user_id)
    except metrics.MetricNotFound as e:
        raise InvalidPlaybook(str(e)) from e

    was = str(block.import_class)
    block.config = {**dict(block.config or {}), "metric_id": metric_id,
                    "imported_values": list(
                        (block.config or {}).get("rows") or [])[:200]}
    block.import_class = "MAPPED_GOVERNED_METRIC"
    block.snapshot_id = None
    session.flush()
    service.record(
        session, entity_type="block", action="mapped", pack=pack,
        entity_id=int(block.id), entity_ref=str(block.block_type),
        changes={"import_class": [was, block.import_class],
                 "metric_id": [None, metric_id]},
        narrative=(f"Mapped to {metric_id}. The values from the uploaded file "
                   "are kept alongside so the two can be compared, and the "
                   "pack will show CreditProbe's figure once it is "
                   "generated."),
        grant=grant)
    return {"block_id": int(block.id), "metric_id": metric_id,
            "import_class": block.import_class,
            "note": ("Generate the pack to calculate it. Until then this "
                     "block has no figure — the imported values are kept as "
                     "a comparison, not as the answer.")}


__all__ = [
    "Imported", "MAX_BYTES", "MAX_RATIO", "MAX_UNPACKED", "MEDIA_HINTS",
    "READABLE", "SUPPORTING", "Unreadable", "import_pack", "inspect",
    "map_to_metric", "read", "safe_name", "sources",
]
