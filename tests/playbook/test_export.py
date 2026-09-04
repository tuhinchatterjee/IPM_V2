"""Four real files, produced from a real pack, and opened again to check them.

Nothing here asserts that a function returned successfully. Every format is
rendered to bytes and then READ BACK with the library that consumes it — the
PDF by its header and length, the Word file and the deck by unzipping them, the
workbook by loading it in openpyxl and looking at the cells. A test that only
checks an export did not raise is a test that passes on an empty file.
"""

from __future__ import annotations

import io

import pytest

from backend.playbook import access, export, generation
from backend.playbook import service as pb

pytestmark = pytest.mark.usefixtures("session")


@pytest.fixture
def calculated(session, pack, actors):
    """A pack with figures in it, some commentary, and a finding or two."""
    generation.generate(session, pack["id"], actors["owner"])
    whole = pb.pack(session, pack["id"], actors["owner"])
    for section in whole["sections"]:
        for block in section["blocks"]:
            if block["block_type"] == "NARRATIVE":
                pb.update_block(
                    session, block["id"], actors["owner"],
                    body="The book performed broadly in line with the "
                         "previous period.")
    return session.get(_pack_model(), int(pack["id"]))


def _pack_model():
    from backend.models.playbook import PlaybookPack

    return PlaybookPack


# =============================================================== the formats


def test_the_pdf_is_a_pdf_with_the_pack_in_it(session, calculated, actors):
    data, media = export.render(session, calculated, "pdf")
    assert media == "application/pdf"
    assert data.startswith(b"%PDF-"), "a real PDF, not an error page"
    assert data.rstrip().endswith(b"%%EOF")
    assert len(data) > 5_000, "a pack with two sections is not a blank page"


def test_the_word_file_opens_and_carries_the_sections(session, calculated,
                                                      actors):
    data, media = export.render(session, calculated, "docx")
    assert "wordprocessingml" in media

    from docx import Document

    document = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in document.paragraphs)
    assert calculated.name in text
    assert "Portfolio performance" in text
    assert "Basis of preparation" in text


def test_the_deck_opens_and_has_a_slide_per_section(session, calculated,
                                                    actors):
    data, media = export.render(session, calculated, "pptx")
    assert "presentationml" in media

    from pptx import Presentation

    deck = Presentation(io.BytesIO(data))
    said = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                said.append(shape.text_frame.text)
    text = "\n".join(said)

    assert calculated.name in text
    assert "Portfolio performance" in text
    assert "Basis of preparation" not in text, (
        "the working belongs in the pack; a slide of formula hashes is a "
        "slide nobody reads")


def test_the_workbook_carries_the_working_behind_every_figure(session,
                                                              calculated,
                                                              actors):
    data, media = export.render(session, calculated, "xlsx")
    assert "spreadsheetml" in media

    from openpyxl import load_workbook

    book = load_workbook(io.BytesIO(data))
    assert set(book.sheetnames) >= {"PACK", "FIGURES", "FINDINGS"}

    figures = book["FIGURES"]
    headers = [c.value for c in figures[1]]
    for wanted in ("Metric", "Period", "Value", "Availability",
                   "Formula hash", "Run id", "Numerator", "Denominator"):
        assert wanted in headers, wanted

    rows = list(figures.iter_rows(min_row=2, values_only=True))
    assert rows, "the pack's figures are in the workbook"
    metrics = {r[headers.index("Metric")] for r in rows}
    assert "retail.default_rate" in metrics


def test_an_unknown_format_is_refused_by_name(session, calculated):
    with pytest.raises(pb.InvalidPlaybook) as e:
        export.render(session, calculated, "csv")
    assert "not a format" in str(e.value)
    assert "pptx" in str(e.value), "the refusal lists what IS available"


# ========================================================= what it must not do


def test_an_export_shows_the_frozen_figures_and_does_not_recalculate(
        session, calculated, actors):
    """The PDF sent in the morning and the screen opened in the afternoon
    show the same numbers, because both read the same snapshots."""
    built = export.document(session, calculated)
    on_screen = pb.pack(session, int(calculated.id), actors["owner"])
    shown = on_screen["sections"][0]["blocks"][0]["figure"]["display_value"]

    table = next(s["table"] for s in built["sections"]
                 if s["title"] == "Portfolio performance")
    assert any(shown in str(cell) for row in table["rows"] for cell in row), (
        f"the pack shows {shown} and the export must show the same string")


def test_an_unaccepted_ai_draft_never_reaches_a_document(session, calculated,
                                                         actors):
    """A DRAFT export is legitimate; carrying words nobody signed is not."""
    whole = pb.pack(session, int(calculated.id), actors["owner"])
    pb.create_block(
        session, whole["sections"][0]["id"], actors["owner"],
        block_type="AI_NARRATIVE",
        body="A machine wrote this and nobody has read it.", source="AI")

    built = export.document(session, calculated)
    text = " ".join(str(s.get("narrative") or "") for s in built["sections"])
    assert "nobody has read it" not in text


def test_a_figure_with_no_value_is_listed_rather_than_omitted(session, pack,
                                                              actors):
    """A pack that simply drops the absent figures looks complete."""
    pb.update_pack(session, pack["id"], actors["owner"], period="2025-07")
    generation.generate(session, pack["id"], actors["owner"])

    from backend.models.playbook import PlaybookPack

    row = session.get(PlaybookPack, int(pack["id"]))
    built = export.document(session, row)
    basis = next(s for s in built["sections"]
                 if s["title"] == "Basis of preparation")
    findings = basis.get("findings") or []
    if not findings:
        pytest.skip("every figure has a value in this lake")
    assert any("has no value" in f["text"] for f in findings)
    assert any("NOT_MATURED" in f["text"] or "NO_DATA" in f["text"]
               or "PERIOD_MISSING" in f["text"] for f in findings)


def test_a_draft_export_says_it_is_a_draft(session, calculated):
    built = export.document(session, calculated)
    assert "DRAFT" in built["purpose"]
    assert "not an approved committee pack" in built["purpose"]


def test_somebody_who_cannot_read_the_pack_cannot_export_it(session, pack,
                                                            actors):
    with pytest.raises(access.PackNotFound):
        export.export(session, pack["id"], actors["outsider"], fmt="pdf")


# ======================================================= Excel and filenames


def test_a_cell_excel_would_execute_is_written_as_text():
    """The formula-injection defence, exercised on the shapes that matter."""
    for dangerous in ("=1+1", "+SUM(A1:A9)", "-0.4pp on the quarter",
                      "@SUM(A1)", "=cmd|'/c calc'!A0"):
        assert export.safe_cell(dangerous).startswith("'"), dangerous


def test_a_legitimate_number_is_not_quoted():
    """Quoting a figure turns it into text nobody can sum."""
    for safe in (6.88, 0, -4, True, None, 19_000):
        assert export.safe_cell(safe) == safe, safe


def test_ordinary_text_passes_through_unchanged():
    for safe in ("Retail default rate", "Q1 2026 review", "6.88%"):
        assert export.safe_cell(safe) == safe


def test_the_defence_survives_into_the_real_workbook(session, calculated,
                                                     actors):
    """Not just the helper — the bytes openpyxl actually wrote."""
    whole = pb.pack(session, int(calculated.id), actors["owner"])
    pb.update_section(session, whole["sections"][0]["id"], actors["owner"],
                      title="=HYPERLINK(\"http://evil\",\"click\")")

    from backend.models.playbook import PlaybookFinding

    session.add(PlaybookFinding(
        pack_id=int(calculated.id), finding_type="DATA_QUALITY",
        severity="LOW", title="-1234 was written to the ledger",
        factual_basis="=1+1", fingerprint="excel-injection-test",
        status="OPEN"))
    session.flush()

    data, _ = export.render(session, calculated, "xlsx")

    from openpyxl import load_workbook

    book = load_workbook(io.BytesIO(data))
    for name in book.sheetnames:
        for row in book[name].iter_rows(values_only=True):
            for cell in row:
                if isinstance(cell, str):
                    assert not cell.startswith(("=", "+", "@")), (
                        f"{name}: {cell!r} would be executed by Excel")
                    assert not (cell.startswith("-")
                                and any(c.isalpha() for c in cell[:12])), cell


def test_a_filename_is_a_name_and_not_a_path(session, calculated):
    calculated.name = "../../etc/passwd"
    session.flush()
    name = export.safe_filename(calculated, "pdf")
    assert "/" not in name and "\\" not in name
    assert ".." not in name
    assert name.endswith(".pdf")


def test_a_title_with_angle_brackets_survives_as_the_title(session, calculated,
                                                           actors):
    """“<Finance> Review” is a legitimate section title.

    Escaped for the document rather than stripped, so the title still reads
    the way the person wrote it — and the PDF build, which parses its strings
    as mini-HTML, does not fall over on it.
    """
    whole = pb.pack(session, int(calculated.id), actors["owner"])
    pb.update_section(session, whole["sections"][0]["id"], actors["owner"],
                      purpose="<Finance> Review & sign-off")

    built = export.document(session, calculated)
    said = next(s for s in built["sections"]
                if "Review" in str(s.get("narrative") or ""))
    assert "&lt;Finance&gt;" in said["narrative"]
    assert "&amp;" in said["narrative"]

    data, _ = export.render(session, calculated, "pdf")
    assert data.startswith(b"%PDF-"), (
        "an angle bracket in a title must not break the PDF build")


# ============================================================== the record


def test_every_export_is_recorded_with_what_was_in_it(session, calculated,
                                                      actors):
    from sqlalchemy import select

    from backend.models.platform import ExportRecord

    outcome = export.export(session, int(calculated.id), actors["owner"],
                            fmt="pdf")
    assert outcome["size"] > 0
    assert len(outcome["checksum"]) == 64

    row = session.execute(
        select(ExportRecord)
        .where(ExportRecord.object_type == "playbook_pack",
               ExportRecord.object_id == str(calculated.id))
        .order_by(ExportRecord.id.desc())).scalars().first()
    assert row is not None
    assert row.content_hash == outcome["checksum"]
    assert row.size_bytes == outcome["size"]
    assert row.user_id == actors["owner"].user_id
    assert "committee_" in row.authorization

    # The figures the file contained, named, so the question "which figures
    # did that file hold" is answerable without regenerating it — which would
    # answer with today's numbers anyway.
    assert row.detail["figures"], row.detail
    assert row.detail["pack_version"] == int(calculated.version)
    assert any(f["metric_id"] == "retail.default_rate"
               for f in row.detail["figures"])


def test_the_export_leaves_a_line_in_the_packs_own_history(session, calculated,
                                                           actors):
    from sqlalchemy import select

    from backend.models.playbook import PlaybookEvent

    export.export(session, int(calculated.id), actors["owner"], fmt="xlsx")
    event = session.execute(
        select(PlaybookEvent).where(
            PlaybookEvent.pack_id == calculated.id,
            PlaybookEvent.action == "exported")
        .order_by(PlaybookEvent.id.desc())).scalars().first()
    assert event is not None
    assert event.author_id == actors["owner"].user_id
    assert "Evidence workbook" in event.narrative


def test_all_four_formats_are_offered_with_what_each_is_for():
    offered = export.formats()
    assert {f["format"] for f in offered} == set(export.FORMATS)
    assert all(f["purpose"] for f in offered)


# =========================================== what an uploaded table must do


@pytest.fixture
def with_imported_table(session, calculated, actors):
    """The pack, plus one table lifted out of somebody's spreadsheet.

    Built through the service the way `import_` builds one, so the block is
    the same shape the importer produces: an UNMAPPED_TABLE carrying its own
    rows, with no snapshot and no metric behind it.
    """
    whole = pb.pack(session, int(calculated.id), actors["owner"])
    section = pb.create_section(session, int(calculated.id), actors["owner"],
                                title="From the branch pack")
    pb.create_block(
        session, int(section["id"]), actors["owner"],
        block_type="TABLE", title="Branch submission",
        import_class="UNMAPPED_TABLE",
        config={
            "imported": True,
            "columns": ["Measure", "Value", "Note"],
            "rows": [
                ["Jeddah SME approval rate", "41.2%", "From the branch pack"],
                ["<Finance> Review", "", "A blank cell, which is ordinary"],
            ],
        })
    assert whole is not None
    return calculated


def test_an_imported_table_reaches_the_document(session, with_imported_table):
    """A table somebody uploaded is content, and content goes in the pack.

    It has no snapshot and no metric — CreditProbe did not calculate it — and
    that is exactly why it must still be carried. Dropping it silently means a
    person puts a table in the pack, sees it on screen, and then cannot find
    it in the file the committee reads.
    """
    built = export.document(session, with_imported_table)
    tables = [s.get("table") for s in built["sections"] if s.get("table")]
    rows = [cell for t in tables for row in t["rows"] for cell in row]
    assert any("Jeddah SME approval rate" in c for c in rows), (
        "the imported rows never reached the document")


def test_an_imported_table_says_where_it_came_from(session,
                                                   with_imported_table):
    """A reader must never mistake it for a CreditProbe calculation."""
    built = export.document(session, with_imported_table)
    titles = [str(s["table"].get("title") or "") for s in built["sections"]
              if s.get("table")]
    assert any("uploaded document" in t for t in titles), titles


def test_a_blank_cell_does_not_break_the_deck(session, with_imported_table):
    """An empty cell in an uploaded table is ordinary, not an exception.

    Setting a PowerPoint cell's text to "" leaves the paragraph with no runs,
    so reaching for `runs[0]` to set the font raises and the whole deck fails.
    A blank cell must not be able to take an export down.
    """
    data, media = export.render(session, with_imported_table, "pptx")
    assert data[:2] == b"PK", "the deck did not build"
    assert "presentation" in media


def test_word_and_the_deck_get_the_characters_the_person_typed(
        session, with_imported_table, actors):
    """Only the PDF parses its input as markup. The others must not see it.

    `document()` escapes because reportlab reads mini-HTML. Word and
    PowerPoint write text literally, so "<Finance> Review" must arrive as
    that and not as "&lt;Finance&gt; Review" — which is what a client would
    otherwise read on the page. Read back with the libraries that consume
    the files rather than out of the raw XML, where the writer's own
    escaping makes the two indistinguishable.
    """
    import docx
    import pptx

    document = docx.Document(io.BytesIO(
        export.render(session, with_imported_table, "docx")[0]))
    words = [c.text for t in document.tables for row in t.rows
             for c in row.cells]
    words += [p.text for p in document.paragraphs]
    assert any("<Finance> Review" in w for w in words), (
        "the characters the person typed did not reach Word")
    assert not any("&lt;" in w or "&amp;" in w for w in words), (
        "Word carried HTML entities into what a person reads")

    deck = pptx.Presentation(io.BytesIO(
        export.render(session, with_imported_table, "pptx")[0]))
    said = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                said.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                said += [c.text for row in shape.table.rows for c in row.cells]
    assert not any("&lt;" in s or "&amp;" in s for s in said), (
        "the deck carried HTML entities into what a person reads")


def test_the_workbook_records_what_a_document_brought_in(session,
                                                        with_imported_table):
    """The evidence file lists the imported rows separately from the figures.

    They are not figures and must never be listed as though they were, but a
    reader checking the pack still has to be able to see what came from a
    file rather than from a calculation.
    """
    import openpyxl

    data, _ = export.render(session, with_imported_table, "xlsx")
    book = openpyxl.load_workbook(io.BytesIO(data))
    assert "IMPORTED" in book.sheetnames, book.sheetnames
    values = [c.value for row in book["IMPORTED"].iter_rows() for c in row]
    assert any(isinstance(v, str) and "Jeddah SME approval rate" in v
               for v in values)
    assert not any(isinstance(v, str) and v[:1] in "=+-@" for v in values), (
        "an imported cell reached the workbook able to start a formula")
